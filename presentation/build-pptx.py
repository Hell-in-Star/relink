# -*- coding: utf-8 -*-
"""
Збирає presentation/relink-presentation.pptx "з нуля" через python-pptx,
за мотивами затвердженого presentation/relink-presentation.html (9 слайдів).

Картинкою вставлені лише: фонові підложки (pptx-assets/bg-*.png),
складна графіка слайда 1 (дуга + коло-схема), маленька лого-мітка
в центрі кола, логотипи ReLink і дрібні іконки-гліфи (усі — pptx-assets/*.png
+ assets/source/лого/png/*.png). Решта — редаговані фігури й текстові поля.

Запуск: python build-pptx.py   (з теки presentation/, після
`node pptx-assets/render-assets.js` і `pip install python-pptx pillow`)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_AUTO_SIZE
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "pptx-assets")
LOGO_DIR = os.path.join(HERE, "..", "assets", "source", "лого", "png")
OUT_PATH = os.path.join(HERE, "relink-presentation.pptx")

# ---------------------------------------------------------------- кольори ---
AMBER = RGBColor(0xDE, 0xA1, 0x02)
GRAPHITE = RGBColor(0x39, 0x40, 0x50)
GREY = RGBColor(0x70, 0x77, 0x87)
LIGHT = RGBColor(0xEB, 0xEA, 0xE8)
DIM = RGBColor(0xA8, 0xAE, 0xBC)
BG_LIGHT = RGBColor(0xEF, 0xEE, 0xEC)
BG_DARK = RGBColor(0x39, 0x40, 0x50)
LINE_GREY = RGBColor(0x70, 0x77, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Geologica Medium"      # заголовки/акценти, Medium
FONT_HEAD_REG = "Geologica"          # Geologica Regular (400)
FONT_BODY_L = "Onest Light"          # основний текст, Light
FONT_BODY_R = "Onest Regular"        # основний текст, Regular
FONT_BODY_M = "Onest Medium"

PX = 96.0  # 96 css-px = 1 inch


def IN(px):
    return Inches(px / PX)


# ============================================================== helpers ===

def new_presentation():
    prs = Presentation()
    prs.slide_width = IN(1280)
    prs.slide_height = IN(720)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # порожній макет
    return prs.slides.add_slide(layout)


def set_alpha(color_format, pct):
    """Додає прозорість (0-100) до вже встановленого RGB-кольору заливки/лінії."""
    srgb = color_format._xFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    for child in list(srgb):
        if child.tag == qn('a:alpha'):
            srgb.remove(child)
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(pct * 1000)))


def add_bg_image(slide, filename):
    slide.shapes.add_picture(os.path.join(ASSETS, filename), 0, 0, width=IN(1280), height=IN(720))


def add_image(slide, path_or_name, x, y, w=None, h=None, in_assets=True):
    p = os.path.join(ASSETS, path_or_name) if in_assets else path_or_name
    kwargs = {}
    if w is not None:
        kwargs['width'] = IN(w)
    if h is not None:
        kwargs['height'] = IN(h)
    return slide.shapes.add_picture(p, IN(x), IN(y), **kwargs)


def add_icon(slide, name, cx, cy, size):
    """Іконка-гліф із pptx-assets/icon-<name>.png, центрована в точці (cx,cy), квадрат size."""
    return add_image(slide, f"icon-{name}.png", cx - size / 2, cy - size / 2, size, size)


def logo_png(variant):
    return os.path.join(LOGO_DIR, f"brand-horizontal-{variant}.png")


LOGO_ASPECT = 766.0 / 228.0  # viewBox вихідного SVG


def add_logo(slide, variant, x, y, height):
    width = height * LOGO_ASPECT
    return slide.shapes.add_picture(logo_png(variant), IN(x), IN(y), width=IN(width), height=IN(height))


def radius_adj(radius_px, w_px, h_px):
    return max(0.0, min(0.5, radius_px / min(w_px, h_px)))


def add_rect(slide, x, y, w, h, fill=None, radius=0, line=None, line_w=1.0,
             line_alpha=None, fill_alpha=None, shadow=False, dash=None, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, IN(x), IN(y), IN(w), IN(h))
    if name:
        shp.name = name
    if radius > 0:
        try:
            shp.adjustments[0] = radius_adj(radius, w, h)
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if fill_alpha is not None:
            set_alpha(shp.fill.fore_color, fill_alpha)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if line_alpha is not None:
            set_alpha(shp.line.color, line_alpha)
        if dash:
            ln = shp.line._get_or_add_ln()
            d = etree.SubElement(ln, qn('a:prstDash'))
            d.set('val', dash)
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def _soft_shadow(shp, blur=10, dist=2, alpha=88, direction=5400000):
    spPr = shp.fill._xPr if hasattr(shp.fill, '_xPr') else shp._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(effectLst, qn('a:outerShdw'))
    sh.set('blurRad', str(Emu(Pt(blur))))
    sh.set('dist', str(Emu(Pt(dist))))
    sh.set('dir', str(direction))
    sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr'))
    clr.set('val', '394050')
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int((100 - alpha) * 1000)))


def add_oval(slide, cx, cy, d, fill=WHITE, shadow=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx - d / 2), IN(cy - d / 2), IN(d), IN(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def add_line(slide, x1, y1, x2, y2, color=LINE_GREY, width_pt=1.0, alpha=None, dash=None, name=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
    if name:
        ln.name = name
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    if alpha is not None:
        set_alpha(ln.line.color, alpha)
    if dash:
        el = ln.line._get_or_add_ln()
        d = etree.SubElement(el, qn('a:prstDash'))
        d.set('val', dash)
    ln.shadow.inherit = False
    return ln


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=0, wrap=True, name=None):
    """
    runs: список параграфів; кожен параграф — список ранів
          (text, font_name, size_pt, color, bold)
          Якщо runs — плаский список ранів (не список списків), трактується як 1 параграф.
    """
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if runs and isinstance(runs[0], tuple):
        runs = [runs]

    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for run_spec in para_runs:
            text, font_name, size, color, bold = run_spec[:5]
            caps = run_spec[5] if len(run_spec) > 5 else False
            letter_spacing = run_spec[6] if len(run_spec) > 6 else None
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            rPr = r._r.get_or_add_rPr()
            if caps:
                rPr.set('cap', 'all')
            if letter_spacing is not None:
                rPr.set('spc', str(int(letter_spacing * 100)))
    return box


def simple_text(slide, x, y, w, h, text, font_name, size, color, bold=False,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True,
                 name=None, caps=False, letter_spacing=None):
    return add_text(slide, x, y, w, h, [(text, font_name, size, color, bold, caps, letter_spacing)],
                     align=align, anchor=anchor, line_spacing=line_spacing, wrap=wrap, name=name)


def add_pill(slide, x, y, w, h, text, border=GREY, text_color=None, fill=None,
             size=10, font_name=FONT_BODY_R, bold=False, name=None, line_w=1.0, caps=True,
             letter_spacing=0.6):
    shp = add_rect(slide, x, y, w, h, fill=fill, radius=h / 2, line=border, line_w=line_w)
    if name:
        shp.name = name
    tc = text_color or border
    simple_text(slide, x, y, w, h, text, font_name, size, tc, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False,
                 caps=caps, letter_spacing=letter_spacing)
    return shp


def pxpt(px):
    """CSS px font-size -> PowerPoint pt (96px/in, 72pt/in)."""
    return px * 0.75


def set_name(shape, name):
    shape.name = name
    return shape


def footer(slide, variant='dark', page_num='1 / 9', height=52):
    """variant: 'dark' -> лого-dark (для світлого низу); 'light' -> лого-light (для темного низу)."""
    logo_variant = 'light' if variant == 'light' else 'dark'
    add_logo(slide, logo_variant, 40, 720 - height + (height - 22) / 2, 22)
    simple_text(slide, 900, 720 - height, 340, height, page_num, FONT_HEAD, pxpt(15), AMBER,
                bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False,
                name="footer-page-num")


# ============================================================ СЛАЙД 1 =====

def build_slide1(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-slide1.png")
    add_image(slide, "slide1-graphic.png", 0, 0, 1280, 720)

    COL_RIGHT_X = 1280 * 0.72  # 921.6

    # top-strip: лінія-роздільник + лого
    add_line(slide, 0, 56, COL_RIGHT_X, 56, color=LINE_GREY, width_pt=0.75, alpha=60)
    add_logo(slide, 'dark', 32, 17, 22)

    # ліва колонка: заголовок-сходинка (bottom-anchored стек, як у HTML)
    simple_text(slide, 40, 455, 320, 50, "Знаходимо", FONT_HEAD, pxpt(44), GRAPHITE,
                wrap=False, name="s1-heading-r1")
    simple_text(slide, 40, 505, 320, 42, "рішення,", FONT_HEAD_REG, pxpt(34), GRAPHITE,
                wrap=False, name="s1-heading-r2")
    simple_text(slide, 40, 546, 320, 50, "що працюють", FONT_HEAD, pxpt(44), AMBER,
                wrap=False, name="s1-heading-r3")

    # пігулки сегментів
    seg_y, seg_h = 620, 27
    seg = [
        (40, 76, "Defense", False), (124, 14, "·", True),
        (146, 99, "Production", False), (253, 14, "·", True),
        (275, 61, "Sales", False), (344, 14, "·", True),
        (366, 108, "Back office", False),
    ]
    for i, (x, w, text, is_dot) in enumerate(seg):
        if is_dot:
            simple_text(slide, x, seg_y, w, seg_h, text, FONT_BODY_R, pxpt(14), AMBER,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        else:
            add_pill(slide, x, seg_y, w, seg_h, text, border=AMBER, text_color=AMBER,
                     size=pxpt(12), font_name=FONT_BODY_M, name=f"s1-seg-pill-{text}")

    # виноски кола-схеми (текст окремо від графіки — редаговані)
    simple_text(slide, 364, 120, 260, 20, "Виробничі", FONT_BODY_R, pxpt(14), GRAPHITE,
                caps=True, letter_spacing=1.4, wrap=False, name="s1-leader-1")
    simple_text(slide, 655, 120, 260, 20, "Інженерно-технічні", FONT_BODY_R, pxpt(14), GRAPHITE,
                align=PP_ALIGN.RIGHT, caps=True, letter_spacing=1.4, wrap=False, name="s1-leader-2")
    simple_text(slide, 364, 487, 260, 20, "Офісні", FONT_BODY_R, pxpt(14), GRAPHITE,
                caps=True, letter_spacing=1.4, wrap=False, name="s1-leader-3")
    simple_text(slide, 655, 487, 260, 20, "Управлінські", FONT_BODY_R, pxpt(14), GRAPHITE,
                align=PP_ALIGN.RIGHT, caps=True, letter_spacing=1.4, wrap=False, name="s1-leader-4")

    # лого-мітка в центрі кола-схеми
    add_image(slide, "center-logo.png", 585, 258, 110, 110)

    # біле коло з іконкою на дузі
    add_oval(slide, 184, 215, 52, fill=WHITE, shadow=True)
    add_icon(slide, "link-amber", 184, 215, 22)

    # права (темна) колонка
    add_logo(slide, 'light', COL_RIGHT_X + 44, 44, 64)
    rb_x = COL_RIGHT_X + 40
    rb_w = 1280 - COL_RIGHT_X - 40 - 18
    simple_text(slide, rb_x, 436, rb_w, 18, "Там, де звичайний найм", FONT_BODY_M, pxpt(13), DIM,
                caps=True, letter_spacing=1.8, wrap=False, name="s1-eyebrow")
    add_pill(slide, rb_x, 467, 118, 30, "не працює", border=None, fill=WHITE, text_color=GRAPHITE,
             size=pxpt(13), font_name=FONT_BODY_M, name="s1-pill-white")
    add_text(slide, rb_x, 515, rb_w, 90, [
        [("власна база ", FONT_BODY_R, pxpt(17), LIGHT, False),
         ("10 000+", FONT_HEAD, pxpt(17), AMBER, False),
         (" кандидатів,", FONT_BODY_R, pxpt(17), LIGHT, False)],
        [("навчальний центр", FONT_BODY_R, pxpt(17), LIGHT, False)],
        [("і підбір під гарантію", FONT_BODY_R, pxpt(17), LIGHT, False)],
    ], line_spacing=1.6, name="s1-right-lines")

    # футер: скрими + роздільник + лого + номер
    add_rect(slide, 0, 664, COL_RIGHT_X, 56, fill=WHITE, fill_alpha=60, name="s1-footer-scrim-light")
    add_rect(slide, COL_RIGHT_X, 664, 1280 - COL_RIGHT_X, 56, fill=GRAPHITE, fill_alpha=60,
             name="s1-footer-scrim-dark")
    add_line(slide, 0, 664, 1280, 664, color=LINE_GREY, width_pt=0.75, alpha=60)
    footer(slide, variant='dark', page_num='1 / 9', height=56)
    return slide


# ============================================================ СЛАЙД 6 =====

CARDS_S6 = [
    ("envelope-amber", "01", "Заявка"),
    ("clipboard-amber", "02", "Бриф"),
    ("search-amber", "03", "Пошук і первинний відбір"),
    ("chat-amber", "04", "Інтерв'ю / скринінг"),
    ("people-amber", "05", "Фіналісти для клієнта"),
    ("checkmark-amber", "06", "Співбесіда"),
    ("document-check-amber", "07", "Вихід кандидата"),
    ("shield-amber", "08", "Гарантійний супровід"),
]


def build_slide6(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-dark-plain.png")

    # заголовок
    simple_text(slide, 80, 40, 1120, 44, "Класичний рекрутинг", FONT_HEAD, pxpt(32), LIGHT,
                caps=True, wrap=False, name="s6-title")
    add_line(slide, 80, 92, 1200, 92, color=AMBER, width_pt=1.5, name="s6-rule")
    add_text(slide, 80, 110, 1120, 50, [
        [("Швидко, прозоро та з гарантією. ", FONT_BODY_M, pxpt(16), AMBER, False),
         ("Офісні, комерційні та управлінські вакансії.", FONT_BODY_L, pxpt(16), LIGHT, False)],
        [("Беремо процес підбору на себе — від уточнення вакансії до виходу кандидата.",
          FONT_BODY_L, pxpt(16), LIGHT, False)],
    ], line_spacing=1.5, name="s6-subtitle")

    # сітка 4x2 карток етапів
    CARD_W, CARD_H, GAP = 265, 190, 20
    for idx, (icon, num, label) in enumerate(CARDS_S6):
        col, row = idx % 4, idx // 4
        cx = 80 + col * (CARD_W + GAP)
        cy = 167 + row * (CARD_H + GAP)
        add_rect(slide, cx, cy, CARD_W, CARD_H, fill=LIGHT, fill_alpha=6, radius=12,
                 name=f"s6-card-{num}")
        add_icon(slide, icon, cx + 28 + 20, cy + 28 + 20, 40)
        simple_text(slide, cx + 28, cy + 28, 209, 44, num, FONT_HEAD, pxpt(44), AMBER,
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, wrap=False,
                    name=f"s6-card-{num}-number")
        simple_text(slide, cx + 28, cy + 90, 209, 72, label, FONT_BODY_R, pxpt(17), LIGHT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, line_spacing=1.3,
                    name=f"s6-card-{num}-name")

    # нижня плашка
    add_rect(slide, 80, 591, 1120, 65, fill=LIGHT, fill_alpha=6, radius=12, name="s6-wide-card")
    add_line(slide, 80, 656, 1200, 656, color=AMBER, width_pt=1.5)
    simple_text(slide, 80, 591, 1120, 65,
                "Ви витрачаєте час лише на співбесіди з кандидатами, які вже пройшли наш відбір.",
                FONT_BODY_R, pxpt(19), LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                name="s6-wide-card-text")

    footer(slide, variant='light', page_num='6 / 9')
    return slide


# ============================================================ СЛАЙД 2 =====

def build_slide2(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-slide2.png")

    COL_X = 1280 * 0.26  # 332.8

    # ліва колонка: заголовок
    add_text(slide, 36, 80, 280, 124, [
        [("Чому", FONT_HEAD, pxpt(48), GRAPHITE, False)],
        [("найм —", FONT_HEAD, pxpt(48), GRAPHITE, False)],
    ], line_spacing=1.08, wrap=False, name="s2-title")
    simple_text(slide, 36, 218, 280, 30, "це критично?", FONT_HEAD_REG, pxpt(24), GREY,
                wrap=False, name="s2-subtitle")
    add_line(slide, COL_X, 0, COL_X, 720, color=LINE_GREY, width_pt=0.75, alpha=60)

    # темна плашка зверху: текст (фон вже в bg-slide2.png)
    tx, tw = COL_X + 64, 1280 - COL_X - 64 - 64
    simple_text(slide, tx, 64, tw, 32, "ReLink — рекрутингова агенція для Defense, виробничих і суміжних компаній.",
                FONT_BODY_R, pxpt(20), LIGHT, line_spacing=1.6, wrap=False, name="s2-p1")
    simple_text(slide, tx, 96, tw, 32, "Ми допомагаємо компаніям, для яких найм персоналу — це не формальність,",
                FONT_BODY_R, pxpt(20), LIGHT, line_spacing=1.6, wrap=False, name="s2-p2")
    simple_text(slide, tx, 142, tw, 32, "а критична умова стабільної роботи й росту бізнесу.",
                FONT_BODY_R, pxpt(20), LIGHT, line_spacing=1.6, wrap=False, name="s2-p3-core")
    simple_text(slide, tx, 192, tw, 40, "Наша задача — зняти з вашої команди операційне навантаження з найму.",
                FONT_HEAD, pxpt(26), AMBER, line_spacing=1.4, wrap=True, name="s2-task")

    add_line(slide, COL_X, 400, 1280, 400, color=LINE_GREY, width_pt=0.75, alpha=60)

    # ланцюжок наслідків (пігулки; стрілки між ними — окремо нижче)
    chain = [
        (463.9, 202, "Незакрита вакансія"),
        (697.9, 202, "Простій виробництва"),
        (943.9, 193, "Втрачені контракти"),
    ]
    for x, w, text in chain:
        add_pill(slide, x, 480, w, 37, text, border=GREY, text_color=GRAPHITE,
                 size=pxpt(15), font_name=FONT_BODY_R, letter_spacing=0.6,
                 name=f"s2-chain-pill-{text}")
    # стрілки-шеврони між пігулками (окремі редаговані лінії-конектори у формі "V")
    for ax in (665.9 + 10, 911.9 + 10):
        cy = 498.5
        add_line(slide, ax - 6, cy - 7, ax + 6, cy, color=AMBER, width_pt=2)
        add_line(slide, ax - 6, cy + 7, ax + 6, cy, color=AMBER, width_pt=2)

    # географія
    add_icon(slide, "globe-amber", 1216 - 175, 655, 16)
    simple_text(slide, 1216 - 160, 646, 160, 18, "Географія — Україна", FONT_BODY_L, pxpt(14), GREY,
                wrap=False, name="s2-geo")

    footer(slide, variant='dark', page_num='2 / 9', height=56)
    return slide


# ============================================================ СЛАЙД 3 =====

def build_cols_header(slide, pill_text, title, pill_style='regular', title_color=LIGHT):
    add_pill(slide, 80, 64, 8 + 13 * len(pill_text), 33, pill_text,
             border=DIM if pill_style == 'regular_dark' else GREY,
             text_color=DIM if pill_style == 'regular_dark' else GREY,
             fill=None, size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0,
             name="header-pill")
    simple_text(slide, 80, 115, 1000, 58, title, FONT_HEAD, pxpt(48), title_color,
                caps=True, wrap=True, name="header-title")


def build_slide3(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-dark-plain.png")
    build_cols_header(slide, "Цільова аудиторія", "З ким працюємо", pill_style='regular_dark', title_color=LIGHT)

    x0, x1 = 80 + 1120 / 3.0, 80 + 2 * 1120 / 3.0
    add_line(slide, x0, 240, x0, 668, color=LINE_GREY, width_pt=0.75, alpha=60)
    add_line(slide, x1, 240, x1, 668, color=LINE_GREY, width_pt=0.75, alpha=60)

    cols = [
        (80, 341.33, "briefcase-graphite", "B2B — клієнти",
         "Керівники HR/HRD та власники або директори компаній сектору Defense (розмір 50–500 осіб) "
         "та МСБ у виробничих, логістичних і технічних сферах (розмір 30–200 осіб)."),
        (485.33, 309.33, "people-graphite", "B2C — кандидати",
         "Робітничі спеціальності (монтажники, комплектувальники), інженерно-технічні фахівці, "
         "офісний та управлінський персонал, ветерани з військовим досвідом."),
        (858.67, 309.33, "veteran-graphite", "Ветеранам",
         "Допомагаємо ветеранам та ветеранкам знайти роботу, де поважають ваш досвід, "
         "відповідальність та навички. Безкоштовно, конфіденційно, крок за кроком."),
    ]
    for i, (x, w, icon, title, text) in enumerate(cols):
        add_oval(slide, x + 32, 240 + 32, 64, fill=WHITE, shadow=True)
        add_icon(slide, icon, x + 32, 240 + 32, 28)
        simple_text(slide, x, 326, w, 30, title, FONT_HEAD, pxpt(22), AMBER,
                    caps=True, wrap=False, name=f"s3-col{i+1}-title")
        simple_text(slide, x, 369, w, 140, text, FONT_BODY_L, pxpt(15), LIGHT,
                    line_spacing=1.6, wrap=True, name=f"s3-col{i+1}-text")

    footer(slide, variant='light', page_num='3 / 9')
    return slide


# ============================================================ СЛАЙД 4 =====

def build_slide4(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-light-plain.png")

    p = add_pill(slide, 80, 64, 130, 36, "Послуги", border=None, fill=None, text_color=AMBER,
                 size=pxpt(16), font_name=FONT_BODY_M, letter_spacing=0.8, name="s4-pill-header")
    p.line.color.rgb = AMBER
    p.line.width = Pt(1)
    set_alpha(p.line.color, 50)
    simple_text(slide, 80, 118, 1000, 58, "Наші послуги", FONT_HEAD, pxpt(48), GRAPHITE,
                caps=True, wrap=True, name="s4-title")

    x0, x1 = 80 + 1120 / 3.0, 80 + 2 * 1120 / 3.0
    add_line(slide, x0, 220, x0, 668, color=LINE_GREY, width_pt=0.75, alpha=60)
    add_line(slide, x1, 220, x1, 668, color=LINE_GREY, width_pt=0.75, alpha=60)

    cols = [
        (80, 341.33, "briefcase-graphite", "Класичний рекрутинг",
         "Офісні, комерційні та управлінські вакансії. Беремо процес підбору на себе — "
         "від уточнення вакансії до виходу кандидата на роботу.", None),
        (485.33, 309.33, "gear-graphite", "Складний інженерно-технічний підбір",
         "Для позицій, де мало резюме і важливий профіль.", None),
        (858.67, 309.33, "link-graphite", "Спеціальні формати співпраці",
         "Масовий підбір, навчальний центр, зовнішній HR-партнер, терміновий підбір.",
         "Детальніше — далі"),
    ]
    for i, (x, w, icon, title, text, note) in enumerate(cols):
        add_oval(slide, x + 32, 220 + 32, 64, fill=WHITE, shadow=True)
        add_icon(slide, icon, x + 32, 220 + 32, 28)
        simple_text(slide, x, 306, w, 52, title, FONT_HEAD, pxpt(19), GRAPHITE,
                    caps=True, wrap=True, line_spacing=1.2, name=f"s4-col{i+1}-title")
        ty = 306 + (24.7 if i == 0 else 49.4) + 12
        simple_text(slide, x, ty, w, 100, text, FONT_BODY_L, pxpt(15), GREY,
                    line_spacing=1.6, wrap=True, name=f"s4-col{i+1}-text")
        if note:
            simple_text(slide, x, ty + 60, w, 26, note, FONT_HEAD, pxpt(18), AMBER,
                        wrap=False, name=f"s4-col{i+1}-note")

    footer(slide, variant='dark', page_num='4 / 9')
    return slide


# ============================================================ СЛАЙД 5 =====

BANDS_S5 = [
    (216, "people-graphite", "Масовий підбір", "При масштабуванні команди."),
    (328, "book-graphite", "Власний навчальний центр", "Підготовка робітничого персоналу."),
    (440, "link-graphite", "Зовнішній HR-партнер", "Для компаній з постійною потребою в наймі."),
    (552, "clock-graphite", "Терміновий підбір", "Коли вакансія горить, до 21 дня від заявки до оферу."),
]


def build_slide5(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-light-plain.png")

    p = add_pill(slide, 80, 64, 230, 36, "Формати співпраці", border=None, fill=None, text_color=AMBER,
                 size=pxpt(16), font_name=FONT_BODY_M, letter_spacing=0.8, name="s5-pill-header")
    p.line.color.rgb = AMBER
    p.line.width = Pt(1)
    set_alpha(p.line.color, 50)
    simple_text(slide, 80, 118, 1100, 58, "Спеціальні формати співпраці", FONT_HEAD, pxpt(48), GRAPHITE,
                caps=True, wrap=True, name="s5-title")

    for i, (top, icon, name, text) in enumerate(BANDS_S5):
        add_line(slide, 80, top, 1200, top, color=LINE_GREY, width_pt=0.75, alpha=60)
        if i == len(BANDS_S5) - 1:
            add_line(slide, 80, top + 112, 1200, top + 112, color=LINE_GREY, width_pt=0.75, alpha=60)
        cy = top + 56
        add_oval(slide, 80 + 28, cy, 56, fill=WHITE, shadow=True)
        add_icon(slide, icon, 80 + 28, cy, 24)
        simple_text(slide, 162, top, 340, 112, name, FONT_HEAD, pxpt(28), GRAPHITE,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, wrap=True, name=f"s5-band{i+1}-name")
        simple_text(slide, 528, top, 672, 112, text, FONT_BODY_L, pxpt(17), GREY,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4, wrap=True, name=f"s5-band{i+1}-text")

    footer(slide, variant='dark', page_num='5 / 9')
    return slide


# ============================================================ СЛАЙД 7 =====

def build_slide7(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-light-plain.png")
    add_rect(slide, 640, 220, 640, 448, fill=GRAPHITE, name="s7-right-panel")
    add_line(slide, 640, 220, 640, 668, color=LINE_GREY, width_pt=0.75, alpha=60)

    add_pill(slide, 80, 64, 190, 33, "Додаткові сервіси", border=GREY, text_color=GREY,
             size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0, name="s7-pill-header")
    simple_text(slide, 80, 115, 900, 116, "Додаткові сервіси за бажанням замовника",
                FONT_HEAD, pxpt(48), GRAPHITE, caps=True, line_spacing=1.1, wrap=True, name="s7-title")

    left_items = [
        "Організація поліграфу", "Перевірка рекомендацій",
        "Допомога з підготовкою Job Offer", "Аудит вакансії або оферу без запуску підбору",
    ]
    add_pill(slide, 80, 260, 150, 33, "Безкоштовні", border=GREY, text_color=GREY,
             size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0, name="s7-left-pill")
    for i, text in enumerate(left_items):
        y = 315 + i * 38
        add_icon(slide, "checkmark-graphite", 80 + 10, y + 10, 20)
        simple_text(slide, 80 + 32, y, 428, 34, text, FONT_BODY_L, pxpt(15), GRAPHITE,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4, wrap=True, name=f"s7-left-item{i+1}")

    right_items = ["Психологічне тестування", "Організація конкурсу або асесменту"]
    add_pill(slide, 696, 260, 110, 33, "Платні", border=LIGHT, text_color=LIGHT,
             fill=None, size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0, name="s7-right-pill")
    for i, text in enumerate(right_items):
        y = 315 + i * 38
        add_icon(slide, "checkmark-light", 696 + 10, y + 10, 20)
        simple_text(slide, 696 + 32, y, 428, 34, text, FONT_BODY_L, pxpt(15), LIGHT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4, wrap=True, name=f"s7-right-item{i+1}")

    footer(slide, variant='dark', page_num='7 / 9')
    return slide


# ============================================================ СЛАЙД 8 =====

def _quad_bezier_points(p0, p1, p2, n=8):
    pts = []
    for i in range(n):
        t = i / (n - 1)
        x = (1 - t) ** 2 * p0[0] + 2 * (1 - t) * t * p1[0] + t ** 2 * p2[0]
        y = (1 - t) ** 2 * p0[1] + 2 * (1 - t) * t * p1[1] + t ** 2 * p2[1]
        pts.append((x, y))
    return pts


STATS_S8 = [("90%", "вакансій закриваються у погоджений термін"),
            ("10 000+", "власна база кандидатів"),
            ("87%", "кандидатів проходять гарантійний термін")]

NODES_S8 = [
    (168, 400, "target-graphite", "Складні вакансії",
     "Досвід закриття інженерно-технічних, технологічних і нестандартних позицій."),
    (468, 370, "lightning-graphite", "Швидкий масовий підбір",
     "Досвід швидкого закриття вакансій маспідбору."),
    (748, 370, "database-graphite", "Власна база кандидатів", "Напрацьована під галузь."),
    (1048, 400, "monitor-graphite", "Доступ до CRM", "Замовник бачить процес підбору в нашій системі."),
]


def build_slide8(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-light-plain.png")

    p = add_pill(slide, 80, 64, 150, 36, "Переваги", border=None, fill=None, text_color=AMBER,
                 size=pxpt(16), font_name=FONT_BODY_M, letter_spacing=0.8, name="s8-pill-header")
    p.line.color.rgb = AMBER
    p.line.width = Pt(1)
    set_alpha(p.line.color, 50)
    simple_text(slide, 80, 118, 700, 58, "Наші переваги", FONT_HEAD, pxpt(48), GRAPHITE,
                caps=True, wrap=False, name="s8-title")

    for i, (num, cap) in enumerate(STATS_S8):
        cx = 80 + 373.33 * i + 373.33 / 2
        simple_text(slide, cx - 170, 190, 340, 52, num, FONT_HEAD, pxpt(52), AMBER,
                    align=PP_ALIGN.CENTER, wrap=False, name=f"s8-stat{i+1}-num")
        simple_text(slide, cx - 170, 250, 340, 30, cap, FONT_BODY_L, pxpt(16), GREY,
                    align=PP_ALIGN.CENTER, line_spacing=1.2, wrap=True, name=f"s8-stat{i+1}-cap")
    add_line(slide, 80 + 373.33, 196, 80 + 373.33, 284, color=LINE_GREY, width_pt=0.75, alpha=60)
    add_line(slide, 80 + 2 * 373.33, 196, 80 + 2 * 373.33, 284, color=LINE_GREY, width_pt=0.75, alpha=60)

    pts = _quad_bezier_points((200, 432), (640, 340), (1080, 432), n=8)
    for i in range(len(pts) - 1):
        add_line(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1],
                 color=LINE_GREY, width_pt=0.75, alpha=60, dash='sysDot')

    for i, (x, y, icon, title, desc) in enumerate(NODES_S8):
        add_oval(slide, x + 32, y + 32, 64, fill=WHITE, shadow=True)
        add_icon(slide, icon, x + 32, y + 32, 28)
        add_text(slide, x - 43, y + 74, 150, 70, [
            [(title, FONT_HEAD, pxpt(15), GRAPHITE, False)],
            [(desc, FONT_BODY_L, pxpt(13), GREY, False)],
        ], align=PP_ALIGN.CENTER, line_spacing=1.35, wrap=True, name=f"s8-node{i+1}")

    footer(slide, variant='dark', page_num='8 / 9')
    return slide


# ============================================================ СЛАЙД 9 =====

def build_slide9(prs):
    slide = blank_slide(prs)
    add_bg_image(slide, "bg-dark-plain.png")
    add_line(slide, 1280 * 0.55, 0, 1280 * 0.55, 720, color=LINE_GREY, width_pt=0.75, alpha=60)

    add_text(slide, 80, 280, 560, 220, [
        [("Готові", FONT_HEAD, pxpt(96), LIGHT, False)],
        [("почати?", FONT_HEAD_REG, pxpt(52), LIGHT, False)],
    ], line_spacing=1.05, wrap=False, name="s9-heading")

    rx, rw = 1280 * 0.55 + 60, 456
    add_rect(slide, rx, 165.5, rw, 230, fill=None, radius=14, line=DIM, line_w=1.0, dash='dash',
             name="s9-card-1")
    add_pill(slide, rx + 30, 165.5 + 26, 150, 33, "Співпраця", border=GREY, text_color=DIM,
             size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0, name="s9-card1-pill")
    fields1 = ["[Ім'я контактної особи]", "[Телефон]", "[E-mail]"]
    for i, f in enumerate(fields1):
        simple_text(slide, rx + 30, 165.5 + 26 + 33 + 18 + i * (24 * 1.3 + 16), rw - 60, 32, f,
                    FONT_HEAD, pxpt(24), LIGHT, wrap=False, name=f"s9-card1-field{i+1}")

    y2 = 165.5 + 230 + 24
    add_rect(slide, rx, y2, rw, 135, fill=None, radius=14, line=DIM, line_w=1.0, dash='dash',
             name="s9-card-2")
    add_pill(slide, rx + 30, y2 + 26, 110, 33, "LinkedIn", border=GREY, text_color=DIM,
             size=pxpt(13), font_name=FONT_BODY_R, letter_spacing=1.0, name="s9-card2-pill")
    simple_text(slide, rx + 30, y2 + 26 + 33 + 18, rw - 60, 32, "[Посилання]",
                FONT_HEAD, pxpt(24), LIGHT, wrap=False, name="s9-card2-field")

    footer(slide, variant='light', page_num='9 / 9')
    return slide


# ============================================================== main ======

def main():
    prs = new_presentation()
    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    build_slide8(prs)
    build_slide9(prs)
    prs.save(OUT_PATH)
    print("OK:", OUT_PATH)


if __name__ == "__main__":
    main()
